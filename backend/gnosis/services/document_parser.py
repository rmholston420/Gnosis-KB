"""
Document parser service.

Supported formats:
  - PDF    — PyMuPDF (fitz)
  - DOCX   — python-docx
  - PPTX   — python-pptx
  - XLSX   — openpyxl
  - Images — pytesseract OCR (PNG, JPG, WEBP, TIFF)
  - URL    — httpx + BeautifulSoup4 scrape

All parsers return a ParsedDocument dataclass with extracted plain text,
detected title, and source metadata.
"""
from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    """Container for a parsed external document."""

    title: str
    text: str
    source: str = ""
    author: str = ""
    page_count: int = 0
    raw_format: str = ""  # pdf | docx | pptx | xlsx | image | url
    metadata: dict = field(default_factory=dict)


# ---------------------------------------------------------------------------
# PDF parser
# ---------------------------------------------------------------------------

def parse_pdf(path: Path) -> ParsedDocument:
    """Extract text from a PDF file using PyMuPDF.

    Args:
        path: Filesystem path to the PDF.

    Returns:
        ParsedDocument with extracted text and page count.
    """
    try:
        import fitz  # type: ignore[import]  # PyMuPDF
    except ImportError:
        raise RuntimeError(
            "PyMuPDF is not installed. Add pymupdf to requirements."
        ) from None

    doc = fitz.open(str(path))
    pages: list[str] = []
    for page in doc:
        pages.append(page.get_text())  # type: ignore[no-untyped-call]

    full_text = "\n\n".join(pages).strip()
    title = Path(path).stem.replace("-", " ").replace("_", " ").title()

    # Try to extract title from PDF metadata
    meta = doc.metadata or {}
    if meta.get("title"):
        title = meta["title"]

    return ParsedDocument(
        title=title,
        text=full_text,
        source=path.name,
        author=meta.get("author", ""),
        page_count=len(doc),
        raw_format="pdf",
        metadata=dict(meta),
    )


# ---------------------------------------------------------------------------
# DOCX parser
# ---------------------------------------------------------------------------

def parse_docx(path: Path) -> ParsedDocument:
    """Extract text from a DOCX file using python-docx.

    Args:
        path: Filesystem path to the DOCX.

    Returns:
        ParsedDocument with extracted paragraph text.
    """
    try:
        from docx import Document  # type: ignore[import]
    except ImportError:
        raise RuntimeError(
            "python-docx is not installed. Add python-docx to requirements."
        ) from None

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)

    # First non-empty paragraph is likely the title
    title = paragraphs[0][:120] if paragraphs else path.stem

    return ParsedDocument(
        title=title,
        text=text,
        source=path.name,
        page_count=len(doc.sections),
        raw_format="docx",
    )


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------

def parse_pptx(path: Path) -> ParsedDocument:
    """Extract text from a PPTX file using python-pptx.

    Args:
        path: Filesystem path to the PPTX.

    Returns:
        ParsedDocument with slide text concatenated.
    """
    try:
        from pptx import Presentation  # type: ignore[import]
    except ImportError:
        raise RuntimeError(
            "python-pptx is not installed. Add python-pptx to requirements."
        ) from None

    prs = Presentation(str(path))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            slides.append(f"## Slide {i}\n" + "\n".join(slide_texts))

    text = "\n\n".join(slides)
    title = path.stem.replace("-", " ").replace("_", " ").title()

    return ParsedDocument(
        title=title,
        text=text,
        source=path.name,
        page_count=len(prs.slides),
        raw_format="pptx",
    )


# ---------------------------------------------------------------------------
# XLSX parser
# ---------------------------------------------------------------------------

def parse_xlsx(path: Path) -> ParsedDocument:
    """Extract data from an XLSX file using openpyxl.

    Converts each sheet to a Markdown table.

    Args:
        path: Filesystem path to the XLSX.

    Returns:
        ParsedDocument with Markdown table representation.
    """
    try:
        import openpyxl  # type: ignore[import]
    except ImportError:
        raise RuntimeError(
            "openpyxl is not installed. Add openpyxl to requirements."
        ) from None

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sections: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))  # type: ignore[arg-type]
        if not rows:
            continue

        # Build Markdown table from first 200 rows
        header = rows[0]
        col_count = len(header)
        md_rows = [
            "| " + " | ".join(str(c) if c is not None else "" for c in header) + " |",
            "| " + " | ".join(["---"] * col_count) + " |",
        ]
        for row in rows[1:201]:  # limit to 200 data rows
            md_rows.append(
                "| " + " | ".join(str(c) if c is not None else "" for c in row) + " |"
            )
        sections.append(f"## Sheet: {sheet_name}\n" + "\n".join(md_rows))

    wb.close()
    text = "\n\n".join(sections)
    title = path.stem.replace("-", " ").replace("_", " ").title()

    return ParsedDocument(
        title=title,
        text=text,
        source=path.name,
        page_count=len(wb.sheetnames),
        raw_format="xlsx",
    )


# ---------------------------------------------------------------------------
# Image / OCR parser
# ---------------------------------------------------------------------------

def parse_image(path: Path) -> ParsedDocument:
    """Extract text from an image using Tesseract OCR.

    Requires tesseract-ocr system package.

    Args:
        path: Filesystem path to the image (PNG / JPG / WEBP / TIFF).

    Returns:
        ParsedDocument with OCR-extracted text.
    """
    try:
        import pytesseract  # type: ignore[import]
        from PIL import Image  # type: ignore[import]
    except ImportError:
        raise RuntimeError(
            "pytesseract or Pillow is not installed."
        ) from None

    img = Image.open(str(path))
    text = pytesseract.image_to_string(img)
    title = path.stem.replace("-", " ").replace("_", " ").title()

    return ParsedDocument(
        title=title,
        text=text.strip(),
        source=path.name,
        page_count=1,
        raw_format="image",
    )


# ---------------------------------------------------------------------------
# URL scraper
# ---------------------------------------------------------------------------

async def parse_url(url: str) -> ParsedDocument:
    """Scrape a URL and extract readable text.

    Uses httpx for the HTTP request and BeautifulSoup4 for HTML parsing.
    Falls back to raw text content if bs4 is unavailable.

    Args:
        url: HTTP/HTTPS URL to scrape.

    Returns:
        ParsedDocument with page title and main text content.
    """
    import httpx

    async with httpx.AsyncClient(timeout=15.0, follow_redirects=True) as client:
        resp = await client.get(url, headers={"User-Agent": "GnosisBot/1.0"})
        resp.raise_for_status()
        html = resp.text

    try:
        from bs4 import BeautifulSoup  # type: ignore[import]

        soup = BeautifulSoup(html, "html.parser")

        # Extract page title
        title_tag = soup.find("title")
        title = title_tag.get_text().strip() if title_tag else url

        # Remove boilerplate elements
        for tag in soup.find_all(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()

        # Try to find main content area
        main = (
            soup.find("main")
            or soup.find("article")
            or soup.find(id="content")
            or soup.find(class_="content")
            or soup.find("body")
        )
        text = main.get_text(separator="\n", strip=True) if main else soup.get_text(separator="\n", strip=True)

    except ImportError:
        logger.warning("beautifulsoup4 not installed — using raw text")
        title = url
        text = html

    return ParsedDocument(
        title=title[:200],
        text=text.strip(),
        source=url,
        raw_format="url",
    )


# ---------------------------------------------------------------------------
# Format dispatcher
# ---------------------------------------------------------------------------

EXTENSION_MAP: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".webp": "image",
    ".tiff": "image",
    ".tif": "image",
}


def detect_format(filename: str) -> Optional[str]:
    """Return the format key for a filename extension, or None if unsupported."""
    ext = Path(filename).suffix.lower()
    return EXTENSION_MAP.get(ext)


def parse_file(path: Path) -> ParsedDocument:
    """Dispatch to the appropriate parser based on file extension.

    Args:
        path: Filesystem path to the file.

    Returns:
        ParsedDocument from the appropriate parser.

    Raises:
        ValueError: If the file extension is not supported.
    """
    fmt = detect_format(path.name)
    if fmt == "pdf":
        return parse_pdf(path)
    if fmt == "docx":
        return parse_docx(path)
    if fmt == "pptx":
        return parse_pptx(path)
    if fmt == "xlsx":
        return parse_xlsx(path)
    if fmt == "image":
        return parse_image(path)
    raise ValueError(f"Unsupported file format: {path.suffix!r}")
